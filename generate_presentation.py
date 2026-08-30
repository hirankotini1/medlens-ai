import os 
from pptx import Presentation 
from pptx .util import Inches ,Pt 
from pptx .enum .text import PP_ALIGN ,MSO_ANCHOR 
from pptx .dml .color import RGBColor 
from pptx .enum .shapes import MSO_SHAPE 

def create_presentation ():
    prs =Presentation ()
    prs .slide_width =Inches (13.333 )
    prs .slide_height =Inches (7.5 )
    blank_layout =prs .slide_layouts [6 ]


    C_BG =RGBColor (11 ,15 ,25 )
    C_CARD =RGBColor (17 ,24 ,39 )
    C_CARD_BORDER =RGBColor (30 ,41 ,59 )
    C_PRIMARY =RGBColor (99 ,102 ,241 )
    C_ACCENT =RGBColor (6 ,182 ,212 )
    C_TEXT_MAIN =RGBColor (248 ,250 ,252 )
    C_TEXT_MUTED =RGBColor (148 ,163 ,184 )
    C_SUCCESS =RGBColor (16 ,185 ,129 )
    C_WARN =RGBColor (245 ,158 ,11 )

    def set_slide_background (slide ):
        bg =slide .shapes .add_shape (MSO_SHAPE .RECTANGLE ,0 ,0 ,Inches (13.333 ),Inches (7.5 ))
        bg .fill .solid ()
        bg .fill .fore_color .rgb =C_BG 
        bg .line .fill .background ()
        return bg 

    def add_header (slide ,title_text ,category_tag ="NEXUS PATHOLOGY"):

        tx_tag =slide .shapes .add_textbox (Inches (0.8 ),Inches (0.4 ),Inches (11.7 ),Inches (0.35 ))
        tf_tag =tx_tag .text_frame 
        tf_tag .word_wrap =True 
        p_tag =tf_tag .paragraphs [0 ]
        p_tag .text =category_tag .upper ()
        p_tag .font .size =Pt (10 )
        p_tag .font .bold =True 
        p_tag .font .color .rgb =C_PRIMARY 


        tx_title =slide .shapes .add_textbox (Inches (0.8 ),Inches (0.7 ),Inches (11.7 ),Inches (0.6 ))
        tf_title =tx_title .text_frame 
        tf_title .word_wrap =True 
        p_title =tf_title .paragraphs [0 ]
        p_title .text =title_text 
        p_title .font .size =Pt (22 )
        p_title .font .bold =True 
        p_title .font .color .rgb =C_TEXT_MAIN 

    def add_card (slide ,left ,top ,width ,height ,title =None ):
        card =slide .shapes .add_shape (MSO_SHAPE .ROUNDED_RECTANGLE ,Inches (left ),Inches (top ),Inches (width ),Inches (height ))
        card .fill .solid ()
        card .fill .fore_color .rgb =C_CARD 
        card .line .color .rgb =C_CARD_BORDER 
        card .line .width =Pt (1.5 )

        if title :
            tx =slide .shapes .add_textbox (Inches (left +0.2 ),Inches (top +0.15 ),Inches (width -0.4 ),Inches (0.4 ))
            tf =tx .text_frame 
            p =tf .paragraphs [0 ]
            p .text =title 
            p .font .size =Pt (13 )
            p .font .bold =True 
            p .font .color .rgb =C_ACCENT 
        return card 




    s1 =prs .slides .add_slide (blank_layout )
    set_slide_background (s1 )


    c1 =add_card (s1 ,1.2 ,1.0 ,10.933 ,5.5 )

    tx1 =s1 .shapes .add_textbox (Inches (1.6 ),Inches (1.3 ),Inches (10.133 ),Inches (4.8 ))
    tf1 =tx1 .text_frame 
    tf1 .word_wrap =True 

    p =tf1 .paragraphs [0 ]
    p .text ="🔬 NEXUS PATHOLOGY"
    p .font .size =Pt (36 )
    p .font .bold =True 
    p .font .color .rgb =C_TEXT_MAIN 

    p2 =tf1 .add_paragraph ()
    p2 .text ="Diagnostic Center & Experimental Clinical ML Decision Support Platform"
    p2 .font .size =Pt (18 )
    p2 .font .bold =True 
    p2 .font .color .rgb =C_PRIMARY 
    p2 .space_before =Pt (10 )

    p3 =tf1 .add_paragraph ()
    p3 .text ="A secure, decoupled digital laboratory management system integrating five disease-specific machine learning pipelines with official pathology workflows."
    p3 .font .size =Pt (12 )
    p3 .font .color .rgb =C_TEXT_MUTED 
    p3 .space_before =Pt (14 )

    p4 =tf1 .add_paragraph ()
    p4 .text ="Project Team: [Student Name(s) / Roll Numbers]  |  Academic Year: 2025–2026\nDepartment: Department of Computer Science & Engineering\nInstitution: [College / University Name]"
    p4 .font .size =Pt (11 )
    p4 .font .color .rgb =C_ACCENT 
    p4 .space_before =Pt (28 )

    p5 =tf1 .add_paragraph ()
    p5 .text ="⚠️ Research & Educational Notice: Developed for academic decision support; not an autonomous clinical diagnostic device."
    p5 .font .size =Pt (9.5 )
    p5 .font .color .rgb =C_WARN 
    p5 .space_before =Pt (20 )




    s2 =prs .slides .add_slide (blank_layout )
    set_slide_background (s2 )
    add_header (s2 ,"Problem Statement & Clinical Motivation")

    cards_s2 =[
    ("📄 Fragmented Laboratory Records","Traditional pathology centers often rely on paper slips or isolated spreadsheets, leading to delayed medical follow-ups, lost historical records, and transcription errors.",0.8 ,1.6 ,5.6 ,2.5 ),
    ("🔒 Insecure Patient Report Access","Patients face significant friction accessing diagnostic findings. Unencrypted email attachments and unauthenticated URLs expose sensitive medical data to unauthorized access.",6.8 ,1.6 ,5.6 ,2.5 ),
    ("📊 Unstructured Biochemical Data","Laboratory results are rarely stored with structured parameters, reference intervals, or flags, preventing automated validation and computational trend analysis.",0.8 ,4.4 ,5.6 ,2.4 ),
    ("⚡ AI Disconnect & Diagnostic Overreach","Existing AI healthcare tools often fail to separate official clinical records from probabilistic outputs, creating medical and legal liabilities without auditability.",6.8 ,4.4 ,5.6 ,2.4 )
    ]
    for title ,desc ,l ,t ,w ,h in cards_s2 :
        add_card (s2 ,l ,t ,w ,h ,title )
        tx =s2 .shapes .add_textbox (Inches (l +0.2 ),Inches (t +0.6 ),Inches (w -0.4 ),Inches (h -0.7 ))
        tf =tx .text_frame 
        tf .word_wrap =True 
        p =tf .paragraphs [0 ]
        p .text =desc 
        p .font .size =Pt (11 )
        p .font .color .rgb =C_TEXT_MUTED 




    s3 =prs .slides .add_slide (blank_layout )
    set_slide_background (s3 )
    add_header (s3 ,"Project Objectives")

    add_card (s3 ,0.8 ,1.6 ,11.7 ,5.2 )
    tx3 =s3 .shapes .add_textbox (Inches (1.1 ),Inches (1.9 ),Inches (11.1 ),Inches (4.6 ))
    tf3 =tx3 .text_frame 
    tf3 .word_wrap =True 

    objs =[
    ("1. Digitize Laboratory Workflows","Provide operational interfaces for staff to register patients, enter structured laboratory parameters, and manage report lifecycles (Draft vs. Finalized)."),
    ("2. Secure Patient Access & Privacy","Implement authenticated patient access with cryptographic PBKDF2 hashing, signed session tokens, and strict IDOR prevention."),
    ("3. Strict Report Decoupling","Guarantee that official laboratory reports remain legally immutable and independent from probabilistic ML inferences."),
    ("4. Five Validated Diagnostic Pipelines","Train, validate, and serialize specialized pipelines for Anemia, Dengue, Liver Disease, Thyroid Profile, and Malaria cell microscopy."),
    ("5. Empirical Synthetic Data Investigation","Conduct controlled synthetic data augmentation experiments (+25%, +50%, +100%) to establish evidence-based production model freezing."),
    ("6. Comprehensive Quality & Security Assurance","Achieve 100% verification across 25 automated security, integration, and ML API test scenarios.")
    ]
    for i ,(title ,desc )in enumerate (objs ):
        p =tf3 .paragraphs [0 ]if i ==0 else tf3 .add_paragraph ()
        p .text =f"{title }: {desc }"
        p .font .size =Pt (11.5 )
        p .font .color .rgb =C_TEXT_MAIN 
        p .space_before =Pt (12 )




    s4 =prs .slides .add_slide (blank_layout )
    set_slide_background (s4 )
    add_header (s4 ,"Existing System vs. Proposed Nexus Pathology")

    add_card (s4 ,0.8 ,1.6 ,5.6 ,5.2 ,"⚠️ Existing Conventional System")
    tx_ex =s4 .shapes .add_textbox (Inches (1.0 ),Inches (2.3 ),Inches (5.2 ),Inches (4.2 ))
    tf_ex =tx_ex .text_frame 
    tf_ex .word_wrap =True 
    for item in [
    "• Manual paper slips or isolated flat spreadsheets",
    "• Insecure delivery via unencrypted emails or physical pickup",
    "• Unstructured text lacking biological reference ranges",
    "• No integrated computational decision-support tools",
    "• Casual record overwriting without audit tracking",
    "• Vulnerable to Insecure Direct Object References (IDOR)"
    ]:
        p =tf_ex .add_paragraph ()if tf_ex .paragraphs [0 ].text else tf_ex .paragraphs [0 ]
        p .text =item 
        p .font .size =Pt (11 )
        p .font .color .rgb =C_TEXT_MUTED 
        p .space_before =Pt (10 )

    add_card (s4 ,6.8 ,1.6 ,5.6 ,5.2 ,"✓ Proposed Nexus Pathology Platform")
    tx_pr =s4 .shapes .add_textbox (Inches (7.0 ),Inches (2.3 ),Inches (5.2 ),Inches (4.2 ))
    tf_pr =tx_pr .text_frame 
    tf_pr .word_wrap =True 
    for item in [
    "• Centralized SQLite database with parameterized queries",
    "• Authenticated Patient Portal with cryptographic PBKDF2 PINs",
    "• Structured clinical reporting with dynamic abnormality flags",
    "• 5 specialized ML decision-support pipelines",
    "• Decoupled immutable reports + independent ML audit log",
    "• Robust RBAC, IDOR protection, and 100% automated test coverage"
    ]:
        p =tf_pr .add_paragraph ()if tf_pr .paragraphs [0 ].text else tf_pr .paragraphs [0 ]
        p .text =item 
        p .font .size =Pt (11 )
        p .font .color .rgb =C_SUCCESS 
        p .space_before =Pt (10 )




    s5 =prs .slides .add_slide (blank_layout )
    set_slide_background (s5 )
    add_header (s5 ,"System Architecture & Decoupled Dataflow")

    add_card (s5 ,0.8 ,1.6 ,11.7 ,5.2 )
    tx5 =s5 .shapes .add_textbox (Inches (1.1 ),Inches (1.8 ),Inches (11.1 ),Inches (4.8 ))
    tf5 =tx5 .text_frame 
    tf5 .word_wrap =True 

    arch_ascii =(
    "                                NEXUS PATHOLOGY WEB PLATFORM\n"
    "                                              │\n"
    "                     ┌────────────────────────┴────────────────────────┐\n"
    "                     │                                                 │\n"
    "               ADMIN PORTAL                                      PATIENT PORTAL\n"
    "       (Staff Login / Report Authoring)                   (Patient ID + Security PIN)\n"
    "                     │                                                 │\n"
    "                     └────────────────────────┬────────────────────────┘\n"
    "                                              ▼\n"
    "                                    FastAPI RESTful Backend\n"
    "                               (RBAC & IDOR Verification Middleware)\n"
    "                                              │\n"
    "                     ┌────────────────────────┴────────────────────────┐\n"
    "                     ▼                                                 ▼\n"
    "           SQLite Persistence Layer                           Machine Learning Layer\n"
    "         (patients, lab_reports, users)                     (5 Validated Scikit-Learn Pipelines)\n"
    "                     │                                                 │\n"
    "                     ▼                                                 ▼\n"
    "          Official Pathology Report                          Experimental ML Decision Support\n"
    "          (Immutable Medical Document)                       (Separate `ml_predictions` Audit Log)\n"
    )
    p =tf5 .paragraphs [0 ]
    p .text =arch_ascii 
    p .font .name ="Consolas"
    p .font .size =Pt (9.5 )
    p .font .color .rgb =C_ACCENT 




    s6 =prs .slides .add_slide (blank_layout )
    set_slide_background (s6 )
    add_header (s6 ,"Technology Stack")

    tech_boxes =[
    ("Frontend Layer","• Semantic HTML5 & Vanilla CSS3\n• Modern ES6+ JavaScript (Fetch API)\n• Responsive CSS Grid & Flexbox\n• Dedicated @media print Stylesheet",0.8 ,1.6 ,3.6 ,2.5 ),
    ("Backend & API Layer","• Python 3.12 Runtime\n• FastAPI 0.115+ (ASGI Framework)\n• Uvicorn 0.34+ (Async Web Server)\n• Pydantic v2 (Runtime Data Validation)",4.8 ,1.6 ,3.6 ,2.5 ),
    ("Database & Persistence","• SQLite 3 Relational Engine\n• Python Native sqlite3 Driver\n• Parameterized Queries (? syntax)\n• Online SQLite Backup API",8.8 ,1.6 ,3.6 ,2.5 ),
    ("Machine Learning","• Scikit-Learn 1.6+ (Pipelines & Models)\n• Joblib 1.4+ (Model Serialization)\n• Pandas 2.2+ & NumPy 2.0+\n• 5 Specialized Classifiers",0.8 ,4.3 ,3.6 ,2.5 ),
    ("Image Processing","• OpenCV 4.10+ (cv2 in-memory decoding)\n• Multi-space color conversions (HSV/LAB)\n• Hu Moments & GLCM Texture\n• 354-D Feature Extractor",4.8 ,4.3 ,3.6 ,2.5 ),
    ("Security & Testing","• PBKDF2-HMAC-SHA256 (100k iters)\n• HMAC Signed Session Tokens\n• FastAPI TestClient & unittest\n• 25/25 Automated Tests Passing",8.8 ,4.3 ,3.6 ,2.5 )
    ]
    for title ,desc ,l ,t ,w ,h in tech_boxes :
        add_card (s6 ,l ,t ,w ,h ,title )
        tx =s6 .shapes .add_textbox (Inches (l +0.15 ),Inches (t +0.55 ),Inches (w -0.3 ),Inches (h -0.65 ))
        tf =tx .text_frame 
        tf .word_wrap =True 
        p =tf .paragraphs [0 ]
        p .text =desc 
        p .font .size =Pt (9.5 )
        p .font .color .rgb =C_TEXT_MUTED 




    s7 =prs .slides .add_slide (blank_layout )
    set_slide_background (s7 )
    add_header (s7 ,"Database Design & Decoupled Schema")

    tbls =[
    ("patients Table","• id (PK)\n• patient_id (UK)\n• name, age, gender\n• contact, email\n• access_pin_hash (PBKDF2)\n• created_at",0.8 ,1.6 ,2.7 ,5.2 ),
    ("lab_reports Table","• id (PK)\n• report_id (UK)\n• patient_id (FK)\n• test_category\n• status (Draft/Finalized)\n• lab_technician, remarks\n• report_data (JSON)\n• created_at, updated_at",3.7 ,1.6 ,2.8 ,5.2 ),
    ("ml_predictions Table","• id (PK)\n• patient_id (FK)\n• report_id (FK)\n• disease, prediction\n• confidence, risk_level\n• model_version, model_used\n• input_snapshot (JSON)\n• disclaimer, created_at",6.7 ,1.6 ,2.9 ,5.2 ),
    ("users Table","• id (PK)\n• username (UK)\n• role (admin/staff)\n• password_hash (PBKDF2)\n• created_at\n\n🛡️ Decoupling Rule:\nML inference never mutates lab_reports.",9.8 ,1.6 ,2.7 ,5.2 )
    ]
    for title ,desc ,l ,t ,w ,h in tbls :
        add_card (s7 ,l ,t ,w ,h ,title )
        tx =s7 .shapes .add_textbox (Inches (l +0.15 ),Inches (t +0.55 ),Inches (w -0.3 ),Inches (h -0.65 ))
        tf =tx .text_frame 
        tf .word_wrap =True 
        p =tf .paragraphs [0 ]
        p .text =desc 
        p .font .size =Pt (9.5 )
        p .font .color .rgb =C_TEXT_MUTED 




    s8 =prs .slides .add_slide (blank_layout )
    set_slide_background (s8 )
    add_header (s8 ,"Machine Learning Methodology & Leakage Prevention")

    add_card (s8 ,0.8 ,1.6 ,11.7 ,5.2 )
    tx8 =s8 .shapes .add_textbox (Inches (1.1 ),Inches (1.8 ),Inches (11.1 ),Inches (4.8 ))
    tf8 =tx8 .text_frame 
    tf8 .word_wrap =True 

    steps =[
    ("1. Data Ingestion & Inspection","Structured datasets inspected for class distributions, physiological boundary validity, and missing values."),
    ("2. Leakage-Free Preprocessing","Scalers and encoders fitted strictly on training partitions. Zero target information used in feature transformations."),
    ("3. Stratified 80/20 Partitioning","Guaranteed proportional target class distribution across training and untouched holdout evaluation partitions."),
    ("4. Algorithm Benchmarking & In-Fold CV","5-Fold Cross-Validation with in-fold transformer fitting to evaluate generalization stability."),
    ("5. Cryptographic Duplicate Auditing","SHA-256 duplicate image detection in Malaria dataset identified and purged 25 duplicate images before retraining."),
    ("6. Pipeline Serialization","Trained Scikit-Learn pipelines serialized to models/ for isolated, zero-retraining runtime execution.")
    ]
    for i ,(title ,desc )in enumerate (steps ):
        p =tf8 .paragraphs [0 ]if i ==0 else tf8 .add_paragraph ()
        p .text =f"• {title }: {desc }"
        p .font .size =Pt (11 )
        p .font .color .rgb =C_TEXT_MAIN 
        p .space_before =Pt (10 )




    s9 =prs .slides .add_slide (blank_layout )
    set_slide_background (s9 )
    add_header (s9 ,"Five Specialized Diagnostic ML Pipelines")

    models_info =[
    ("🩸 Anemia (CBC)","Algorithm: Logistic Regression\nInputs: 11 Features (HGB, RBC, PCV, MCV, MCH, MCHC, RDW, TLC, PLT, Age, Sex)\nTarget: Non-Anemic vs. Anemic\nRationale: High linear separability of erythrocyte indices.",0.8 ,1.6 ,5.6 ,2.5 ),
    ("🦟 Dengue Hematology","Algorithm: Random Forest Classifier\nInputs: 8 Features (Platelets, WBC, Differential, RBC, PDW, Hemoglobin, Age, Gender)\nTarget: Negative vs. Positive\nRationale: Captures non-linear platelet-leukocyte interactions.",6.8 ,1.6 ,5.6 ,2.5 ),
    ("🫁 Liver Disease (LFT)","Algorithm: Gradient Boosting Classifier\nInputs: 10 Features (Bilirubin, ALT, AST, ALP, Proteins, Albumin, A/G Ratio, Age, Gender)\nTarget: Non-Liver vs. Liver Disease\nRationale: High sensitivity (95.06% Recall) on complex boundaries.",0.8 ,4.3 ,5.6 ,2.5 ),
    ("🦋 Thyroid Profile","Algorithm: Multinomial Logistic Regression\nInputs: 5 Features (TSH, T4, T3, TSH Response, T3 Resin Uptake)\nTarget: 3 Classes (Normal, Hyperthyroid, Hypothyroid)\nRationale: Softmax probabilities on clustered hormone spectra.",6.8 ,4.3 ,5.6 ,1.2 ),
    ("🔬 Malaria (Microscopy)","Algorithm: Gradient Boosting + 354-D CV Extractor\nInputs: Blood smear cell image (Color stats, HSV/LAB histograms, Hu moments, GLCM)\nTarget: Uninfected vs. Parasitized (97.80% Recall)",6.8 ,5.6 ,5.6 ,1.2 )
    ]
    for title ,desc ,l ,t ,w ,h in models_info :
        add_card (s9 ,l ,t ,w ,h ,title )
        tx =s9 .shapes .add_textbox (Inches (l +0.15 ),Inches (t +0.45 ),Inches (w -0.3 ),Inches (h -0.55 ))
        tf =tx .text_frame 
        tf .word_wrap =True 
        p =tf .paragraphs [0 ]
        p .text =desc 
        p .font .size =Pt (9 )
        p .font .color .rgb =C_TEXT_MUTED 




    s10 =prs .slides .add_slide (blank_layout )
    set_slide_background (s10 )
    add_header (s10 ,"Model Results & Performance Evaluation")

    add_card (s10 ,0.8 ,1.6 ,11.7 ,5.2 )
    tx10 =s10 .shapes .add_textbox (Inches (1.1 ),Inches (1.8 ),Inches (11.1 ),Inches (4.8 ))
    tf10 =tx10 .text_frame 
    tf10 .word_wrap =True 

    res_text =(
    "| Disease / Panel     | Selected Algorithm        | Holdout Accuracy | 5-Fold Cross-Validation | Primary Metric |\n"
    "|---------------------|---------------------------|------------------|-------------------------|----------------|\n"
    "| Anemia (CBC)        | Logistic Regression       | 100.00%          | 95.49% ± 1.64%          | F1: 100.00%    |\n"
    "| Dengue Hematology   | Random Forest Classifier  | 92.93%           | 91.30% ± 2.36%          | Recall: 93.10% |\n"
    "| Liver Disease (LFT) | Gradient Boosting         | 72.81%           | 69.30% ± 2.94%          | Recall: 95.06% |\n"
    "| Thyroid Profile     | Multinomial Logistic Reg  | 100.00%          | 95.81% ± 3.09%          | F1: 100.00%    |\n"
    "| Malaria Microscopy  | Gradient Boosting + CV    | 94.03% (Unseen)  | Strict Deduplicated Set  | Recall: 97.80% |\n\n"
    "Key Engineering & Evaluation Insights:\n"
    "• Malaria Evaluation: Achieved 94.03% accuracy, 93.68% precision, 97.80% recall, and 95.70% F1 on strict unseen images.\n"
    "• Liver Optimization: Prioritized Sensitivity (95.06% Recall) to minimize dangerous False Negatives in hepatic triage.\n"
    "• ⚠️ Academic Clarification: 100% holdout scores reflect curated dataset separability — not clinical diagnostic certainty."
    )
    p =tf10 .paragraphs [0 ]
    p .text =res_text 
    p .font .name ="Consolas"
    p .font .size =Pt (10 )
    p .font .color .rgb =C_TEXT_MAIN 




    s11 =prs .slides .add_slide (blank_layout )
    set_slide_background (s11 )
    add_header (s11 ,"Controlled Synthetic Data Experiment & Findings")

    add_card (s11 ,0.8 ,1.6 ,5.6 ,5.2 ,"🔬 Experimental Design")
    tx_se =s11 .shapes .add_textbox (Inches (1.0 ),Inches (2.2 ),Inches (5.2 ),Inches (4.3 ))
    tf_se =tx_se .text_frame 
    tf_se .word_wrap =True 
    for item in [
    "• Goal: Determine whether synthetic data improves the four tabular models.",
    "• Tested Ratios: +25%, +50%, and +100% synthetic augmentations.",
    "• Strict Quarantine: Real 20% holdout test partition remained untouched.",
    "• In-Fold Generation: Synthesis performed strictly inside training folds during 5-fold CV.",
    "• Domain Clipping: Values constrained to realistic clinical min/max boundaries."
    ]:
        p =tf_se .add_paragraph ()if tf_se .paragraphs [0 ].text else tf_se .paragraphs [0 ]
        p .text =item 
        p .font .size =Pt (10.5 )
        p .font .color .rgb =C_TEXT_MUTED 
        p .space_before =Pt (8 )

    add_card (s11 ,6.8 ,1.6 ,5.6 ,5.2 ,"📊 Results & Production Decision")
    tx_sd =s11 .shapes .add_textbox (Inches (7.0 ),Inches (2.2 ),Inches (5.2 ),Inches (4.3 ))
    tf_sd =tx_sd .text_frame 
    tf_sd .word_wrap =True 
    for item in [
    "• Liver Degradation: Holdout accuracy dropped from 72.81% down to 66.67% with 100% synthetic data.",
    "• Neutral on Linear Separations: Anemia and Thyroid showed zero statistical gain.",
    "• Dengue Degradation: Holdout accuracy dropped from 92.93% to 90.91%.",
    "• FINAL DECISION: Synthetic augmentation did NOT improve real generalization.",
    "• PRODUCTION RULE: Models are frozen strictly on real clinical data. Synthetic experiments remain research-only."
    ]:
        p =tf_sd .add_paragraph ()if tf_sd .paragraphs [0 ].text else tf_sd .paragraphs [0 ]
        p .text =item 
        p .font .size =Pt (10.5 )
        p .font .color .rgb =C_WARN if "FINAL"in item or "PRODUCTION"in item else C_TEXT_MAIN 
        p .space_before =Pt (8 )




    s12 =prs .slides .add_slide (blank_layout )
    set_slide_background (s12 )
    add_header (s12 ,"Cybersecurity & Data Privacy Controls")

    sec_cards =[
    ("🔐 Authentication & Hashing","PBKDF2-HMAC-SHA256 (100,000 iterations + 16-byte random salts) for staff passwords and patient PINs. Plaintext credentials are never stored.",0.8 ,1.6 ,5.6 ,2.5 ),
    ("🛡️ RBAC & IDOR Defense","Signed HMAC Bearer session tokens. Patient A cannot access Patient B's reports or ML history (rejected with HTTP 403 Forbidden).",6.8 ,1.6 ,5.6 ,2.5 ),
    ("💉 SQLi & Input Validation","Parameterized SQL queries (? syntax) neutralize SQL injection. Pydantic schemas enforce physiological boundary validation.",0.8 ,4.3 ,5.6 ,2.5 ),
    ("🩸 Image & Report Integrity","Malaria upload capped at 5MB with OpenCV in-memory validation. Official lab reports are strictly immutable by ML analysis.",6.8 ,4.3 ,5.6 ,2.5 )
    ]
    for title ,desc ,l ,t ,w ,h in sec_cards :
        add_card (s12 ,l ,t ,w ,h ,title )
        tx =s12 .shapes .add_textbox (Inches (l +0.2 ),Inches (t +0.6 ),Inches (w -0.4 ),Inches (h -0.7 ))
        tf =tx .text_frame 
        tf .word_wrap =True 
        p =tf .paragraphs [0 ]
        p .text =desc 
        p .font .size =Pt (10.5 )
        p .font .color .rgb =C_TEXT_MUTED 




    s13 =prs .slides .add_slide (blank_layout )
    set_slide_background (s13 )
    add_header (s13 ,"Clinical Web Application Workflow")

    add_card (s13 ,0.8 ,1.6 ,11.7 ,5.2 )
    tx13 =s13 .shapes .add_textbox (Inches (1.1 ),Inches (1.8 ),Inches (11.1 ),Inches (4.8 ))
    tf13 =tx13 .text_frame 
    tf13 .word_wrap =True 

    flow_text =(
    "ADMIN / LAB STAFF WORKFLOW:\n"
    "  [1. Staff Login] ──> [2. Register Patient] ──> [3. Author Lab Report with Dynamic Panels]\n"
    "           │\n"
    "           ▼\n"
    "  [4. Save as Draft or Finalize Report] ──> [5. Run Experimental ML Decision Support]\n"
    "           │\n"
    "           ▼\n"
    "  [6. ML Inferences Stored in Separate `ml_predictions` Table with Input Snapshots]\n\n"
    "PATIENT WORKFLOW:\n"
    "  [1. Patient Login via ID + PIN] ──> [2. Access Personalized Dashboard (IDOR-Protected)]\n"
    "           │\n"
    "           ▼\n"
    "  [3. View Official Pathology Report Sheet with Biological Reference Ranges & Doctor Remarks]\n"
    "           │\n"
    "           ▼\n"
    "  [4. View Decoupled ML Decision Support Card (Confidence, Risk Level, Model Provenance & Disclaimer)]\n"
    "           │\n"
    "           ▼\n"
    "  [5. Generate Formatted Physical Laboratory Printout (@media print)]"
    )
    p =tf13 .paragraphs [0 ]
    p .text =flow_text 
    p .font .name ="Consolas"
    p .font .size =Pt (10 )
    p .font .color .rgb =C_ACCENT 




    s14 =prs .slides .add_slide (blank_layout )
    set_slide_background (s14 )
    add_header (s14 ,"Testing, Verification & Quality Assurance")

    test_boxes =[
    ("Security Audit Suite","• 15 / 15 Tests Passed (100%)\n• Unauthorized access rejection (401)\n• IDOR cross-patient defense (403)\n• SQL injection payload immunity\n• Fake image & oversized upload rejection\n• File isolation & backup routines",0.8 ,1.6 ,3.6 ,5.2 ),
    ("Pathology Integration Suite","• 10 / 10 Tests Passed (100%)\n• Patient registration & retrieval\n• Report creation & status lifecycle\n• Anemia, Dengue, Liver & Thyroid ML\n• Malaria image inference pipeline\n• Missing value error handling (422)",4.8 ,1.6 ,3.6 ,5.2 ),
    ("Consolidated QA Benchmark","• 25 Distinct End-to-End Scenarios\n• 100% Pass Rate in 0.35s\n• Report Immutability Verified (Bit-for-bit equivalence before & after ML)\n• Zero Model Retraining at Runtime\n• Pydantic Boundary Validation Active",8.8 ,1.6 ,3.6 ,5.2 )
    ]
    for title ,desc ,l ,t ,w ,h in test_boxes :
        add_card (s14 ,l ,t ,w ,h ,title )
        tx =s14 .shapes .add_textbox (Inches (l +0.15 ),Inches (t +0.55 ),Inches (w -0.3 ),Inches (h -0.65 ))
        tf =tx .text_frame 
        tf .word_wrap =True 
        p =tf .paragraphs [0 ]
        p .text =desc 
        p .font .size =Pt (10 )
        p .font .color .rgb =C_SUCCESS if "Passed"in desc or "100%"in desc else C_TEXT_MUTED 
        p .space_before =Pt (4 )




    s15 =prs .slides .add_slide (blank_layout )
    set_slide_background (s15 )
    add_header (s15 ,"Conclusion & Future Scope")

    add_card (s15 ,0.8 ,1.6 ,5.6 ,5.2 ,"📌 Project Conclusion")
    tx_co =s15 .shapes .add_textbox (Inches (1.0 ),Inches (2.2 ),Inches (5.2 ),Inches (4.3 ))
    tf_co =tx_co .text_frame 
    tf_co .word_wrap =True 
    for item in [
    "• Successfully digitized pathology laboratory workflows with accredited reporting standards.",
    "• Integrated 5 specialized ML decision-support pipelines with high diagnostic sensitivity.",
    "• Established a strictly decoupled architecture preserving official report immutability.",
    "• Implemented cryptographic security, RBAC, and IDOR defenses with 100% test pass rate.",
    "• Formulates a responsible, transparent blueprint for AI in medical informatics."
    ]:
        p =tf_co .add_paragraph ()if tf_co .paragraphs [0 ].text else tf_co .paragraphs [0 ]
        p .text =item 
        p .font .size =Pt (10.5 )
        p .font .color .rgb =C_TEXT_MAIN 
        p .space_before =Pt (8 )

    add_card (s15 ,6.8 ,1.6 ,5.6 ,5.2 ,"🚀 Future Scope (Future Work)")
    tx_fu =s15 .shapes .add_textbox (Inches (7.0 ),Inches (2.2 ),Inches (5.2 ),Inches (4.3 ))
    tf_fu =tx_fu .text_frame 
    tf_fu .word_wrap =True 
    for item in [
    "• Explainable AI: Integrate SHAP/LIME feature attribution plots for clinician transparency.",
    "• Hospital Interoperability: Implement HL7 / FHIR v4 standards for EHR communication.",
    "• Whole-Slide Imaging (WSI): Deep learning CNNs/ViTs for gigapixel digital microscopy.",
    "• Enterprise Backend: Migration to clustered PostgreSQL with encrypted storage (TDE).",
    "• Regulatory Compliance: Formal clinical validation and SaMD evaluation (CDSCO / FDA)."
    ]:
        p =tf_fu .add_paragraph ()if tf_fu .paragraphs [0 ].text else tf_fu .paragraphs [0 ]
        p .text =item 
        p .font .size =Pt (10.5 )
        p .font .color .rgb =C_ACCENT 
        p .space_before =Pt (8 )

    os .makedirs ("presentation",exist_ok =True )
    out_path =os .path .abspath ("presentation/Nexus_Pathology_Presentation.pptx")
    prs .save (out_path )
    print (f"Presentation saved successfully to: {out_path }")

if __name__ =="__main__":
    create_presentation ()
